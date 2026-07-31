#!/usr/bin/env python3
"""
PPT 組裝指令碼

將一個目錄中的投影片圖片組裝成 PowerPoint 簡報。
每張圖片會被插入為一頁投影片，充滿整個頁面。
"""

import argparse
import os
import re
import sys
import tempfile
from typing import Dict, List, Optional, Tuple


def dependency_hint() -> str:
    runtime_home = os.path.expanduser(os.environ.get("CODEX_PPT_HOME", "~/.codex-ppt-skill"))
    python = os.path.join(
        runtime_home,
        ".venv",
        "Scripts" if os.name == "nt" else "bin",
        "python.exe" if os.name == "nt" else "python",
    )
    runtime_script = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "codex_ppt_runtime.py",
    )
    return (
        f"請執行: python3 {runtime_script} bootstrap\n"
        f"或直接執行: {python} -m pip install -r "
        f"{os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'requirements.txt')}"
    )


def get_slide_images(ppt_project_dir: str) -> List[str]:
    """
    獲取投影片圖片檔案列表，按頁碼排序

    Args:
        ppt_project_dir: PPT 專案目錄（包含 origin_image 子目錄）

    Returns:
        按順序排列的圖片檔案路徑列表
    """
    # 從 origin_image 子目錄讀取
    origin_image_dir = os.path.join(ppt_project_dir, "origin_image")

    if not os.path.exists(origin_image_dir):
        print(f"錯誤：origin_image 目錄不存在: {origin_image_dir}")
        return []

    print(f"從 origin_image 目錄讀取圖片: {origin_image_dir}")

    # 支援的圖片格式
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp'}

    # 只獲取正式投影片圖片，避免 sample_slide.png、草稿圖或參考圖被誤裝入 PPT。
    slide_name_pattern = re.compile(r"^slide_(\d+)\.(png|jpe?g|gif|bmp)$", re.IGNORECASE)
    image_files = []
    for file in os.listdir(origin_image_dir):
        file_path = os.path.join(origin_image_dir, file)
        if os.path.isfile(file_path):
            ext = os.path.splitext(file)[1].lower()
            if ext in image_extensions and slide_name_pattern.match(file):
                image_files.append(file_path)

    def slide_sort_key(path: str) -> Tuple[int, str]:
        filename = os.path.basename(path)
        stem = os.path.splitext(filename)[0]
        return (int(re.search(r"^slide_(\d+)", stem).group(1)), filename.lower())

    # 按頁碼排序，避免 slide_10 排在 slide_2 前面。
    # 同頁碼多檔案時按檔名穩定排序。
    image_files.sort(key=slide_sort_key)

    return image_files


def load_speaker_notes(ppt_project_dir: str) -> Dict[int, str]:
    """
    從 speech.md 讀取每頁演講備註。

    支援以下標題格式：
    - ## Slide 1: 標題
    - ### Slide 1: 標題
    - ## 第 1 頁：標題
    """
    speech_path = os.path.join(ppt_project_dir, "speech.md")
    if not os.path.exists(speech_path):
        return {}

    with open(speech_path, "r", encoding="utf-8") as f:
        content = f.read()

    notes: Dict[int, str] = {}
    current_slide: Optional[int] = None
    current_lines: List[str] = []
    heading_pattern = re.compile(r"^#{2,4}\s*(?:Slide\s*(\d+)|第\s*(\d+)\s*[頁页])\b.*$", re.IGNORECASE)

    def flush_current() -> None:
        if current_slide is None:
            return
        text = "\n".join(current_lines).strip()
        if text:
            notes[current_slide] = text

    for line in content.splitlines():
        match = heading_pattern.match(line.strip())
        if match:
            flush_current()
            current_slide = int(match.group(1) or match.group(2))
            current_lines = []
            continue

        if current_slide is not None:
            current_lines.append(line)

    flush_current()
    return notes


def compress_image_if_needed(
    image_path: str,
    max_size_mb: float = 2.0,
    quality_step: int = 5
) -> Optional[str]:
    """
    如果圖片超過指定大小，壓縮圖片並返回臨時檔案路徑

    Args:
        image_path: 原始圖片路徑
        max_size_mb: 最大檔案大小（MB）
        quality_step: 每次降低的品質步進

    Returns:
        str: 如果需要壓縮，返回臨時檔案路徑；否則返回 None
    """
    max_size_bytes = max_size_mb * 1024 * 1024

    # 檢查原始檔案大小
    file_size = os.path.getsize(image_path)

    if file_size <= max_size_bytes:
        # 不需要壓縮
        return None

    print(f"  圖片大小 {file_size / 1024 / 1024:.2f}MB，需要壓縮...")

    try:
        from PIL import Image

        # 開啟圖片
        img = Image.open(image_path)

        # 轉換 RGBA 到 RGB（如果需要儲存為 JPEG）
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
            img = background

        # 建立臨時檔案
        temp_fd, temp_path = tempfile.mkstemp(suffix='.jpg')
        os.close(temp_fd)

        # 從高品質開始嘗試壓縮
        quality = 95
        while quality > 20:
            img.save(temp_path, 'JPEG', quality=quality, optimize=True)
            compressed_size = os.path.getsize(temp_path)

            if compressed_size <= max_size_bytes:
                print(f"  壓縮成功: {compressed_size / 1024 / 1024:.2f}MB (品質: {quality})")
                return temp_path

            quality -= quality_step

        # 如果還是太大，嘗試調整尺寸
        print(f"  品質壓縮不足，嘗試縮小尺寸...")
        scale = 0.9
        while scale > 0.3:
            new_width = int(img.width * scale)
            new_height = int(img.height * scale)
            resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

            resized_img.save(temp_path, 'JPEG', quality=85, optimize=True)
            compressed_size = os.path.getsize(temp_path)

            if compressed_size <= max_size_bytes:
                print(f"  壓縮成功: {compressed_size / 1024 / 1024:.2f}MB (縮放: {scale:.0%})")
                return temp_path

            scale -= 0.1

        # 實在壓不下去了，返回最後的結果
        print(f"  警告：無法壓縮到 {max_size_mb}MB 以下，使用最小尺寸版本")
        return temp_path

    except ImportError:
        print("錯誤：未安裝 Pillow 庫")
        print(dependency_hint())
        return None
    except Exception as e:
        print(f"  警告：圖片壓縮失敗: {e}")
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.remove(temp_path)
        return None


def create_presentation(
    image_files: List[str],
    output_path: str,
    aspect_ratio: str = "16:9",
    speaker_notes: Optional[Dict[int, str]] = None,
) -> bool:
    """
    建立 PowerPoint 簡報

    Args:
        image_files: 投影片圖片檔案列表
        output_path: 輸出 PPT 檔案路徑
        aspect_ratio: 投影片寬高比（16:9 或 4:3）

    Returns:
        bool: 成功返回 True，失敗返回 False
    """
    try:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            print("錯誤：未安裝 python-pptx 庫")
            print(dependency_hint())
            return False

        # 建立簡報
        prs = Presentation()

        # 設定投影片尺寸
        if aspect_ratio == "16:9":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
        elif aspect_ratio == "4:3":
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        else:
            print(f"警告：不支援的寬高比 {aspect_ratio}，使用預設值 16:9")
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

        speaker_notes = speaker_notes or {}

        # 新增每一頁投影片
        temp_files_to_cleanup = []
        for i, image_path in enumerate(image_files, 1):
            if not os.path.exists(image_path):
                print(f"警告：圖片檔案不存在: {image_path}")
                continue

            # 壓縮圖片（如果需要）
            compressed_path = compress_image_if_needed(image_path, max_size_mb=2.0)

            # 使用壓縮後的圖片或原圖
            image_to_use = compressed_path if compressed_path else image_path

            # 記錄臨時檔案以便後續清理
            if compressed_path:
                temp_files_to_cleanup.append(compressed_path)

            # 使用空白布局（索引 6）
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 將圖片新增到投影片，充滿整個頁面
            slide.shapes.add_picture(
                image_to_use,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )

            note_text = speaker_notes.get(i)
            if note_text:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.clear()
                notes_frame.text = note_text

            print(f"✓ 已新增第 {i} 頁: {os.path.basename(image_path)}")

        # 確保輸出目錄存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # 儲存簡報
        prs.save(output_path)
        print(f"\n✓ PPT 檔案已儲存: {output_path}")
        print(f"  總頁數: {len(image_files)}")
        if speaker_notes:
            matched_notes = sum(1 for i in range(1, len(image_files) + 1) if speaker_notes.get(i))
            print(f"  已寫入備註: {matched_notes}/{len(image_files)} 頁")

        # 清理臨時檔案
        for temp_file in temp_files_to_cleanup:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            except Exception as e:
                print(f"警告：清理臨時檔案失敗 {temp_file}: {e}")

        return True

    except Exception as e:
        print(f"錯誤：建立 PPT 失敗: {e}")
        import traceback
        traceback.print_exc()

        # 即使失敗也要清理臨時檔案
        if 'temp_files_to_cleanup' in locals():
            for temp_file in temp_files_to_cleanup:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except:
                    pass

        return False


def main():
    parser = argparse.ArgumentParser(
        description='將投影片圖片組裝成 PowerPoint 簡報',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
範例用法:
  # 基本用法
  # 會在 /path/to/base/ 下尋找 MyPresentation/ 資料夾
  # 從 MyPresentation/origin_image/ 讀取圖片
  # 將 PPT 儲存為 MyPresentation/MyPresentation.pptx
  python assemble_ppt.py /path/to/base/ MyPresentation.pptx

  # 指定 4:3 寬高比
  python assemble_ppt.py /path/to/base/ MyPresentation.pptx --aspect-ratio 4:3

  # 只初始化目錄，不生成 PPT
  python assemble_ppt.py /path/to/base/ MyPresentation.pptx --init

資料夾結構要求:
  /path/to/base/MyPresentation/
  ├── origin_image/
  │   ├── slide_01.png
  │   ├── slide_02.png
  │   └── ...
  ├── speech.md (可選，將寫入 PPT 備註)
  └── MyPresentation.pptx (將在此生成)

注意：
  - 圖片檔案必須放在 origin_image 子目錄中
  - 只會讀取 slide_01.png、slide_02.png 這類正式圖片，其他圖片會被忽略
  - 圖片檔案按頁碼排序
  - 建議圖片檔案命名為: slide_01.png, slide_02.png, ...
  - 每張圖片會充滿整個投影片頁面
  - 如果專案目錄下存在 speech.md，會按 Slide N 標題寫入每頁備註
        '''
    )

    parser.add_argument('base_dir', help='基礎目錄（PPT 專案資料夾的父目錄）')
    parser.add_argument('output', help='輸出 PPT 檔名 (.pptx)')
    parser.add_argument('--aspect-ratio', '--ar',
                        choices=['16:9', '4:3'],
                        default='16:9',
                        help='投影片寬高比 (預設: 16:9)')
    parser.add_argument('--init',
                        action='store_true',
                        help='只建立 PPT 專案目錄和 origin_image 子目錄，不生成 PPT')

    args = parser.parse_args()

    # 確保輸出檔案有 .pptx 副檔名
    output_filename = args.output
    if not output_filename.lower().endswith('.pptx'):
        output_filename += '.pptx'

    # 獲取 PPT 名稱（不含副檔名）
    ppt_name = os.path.splitext(os.path.basename(output_filename))[0]

    # 構建 PPT 專案目錄
    ppt_project_dir = os.path.join(args.base_dir, ppt_name)
    origin_image_dir = os.path.join(ppt_project_dir, "origin_image")

    if args.init:
        os.makedirs(origin_image_dir, exist_ok=True)
        print(f"✓ PPT 專案目錄已準備好: {ppt_project_dir}")
        print(f"✓ 投影片圖片目錄已準備好: {origin_image_dir}")
        sys.exit(0)

    if not os.path.exists(ppt_project_dir):
        print(f"錯誤：PPT 專案目錄不存在: {ppt_project_dir}")
        print("如需初始化目錄，請加入 --init 參數")
        sys.exit(1)

    if not os.path.exists(origin_image_dir):
        print(f"錯誤：origin_image 目錄不存在: {origin_image_dir}")
        print("如需初始化目錄，請加入 --init 參數")
        sys.exit(1)

    # 設定輸出路徑
    output_path = os.path.join(ppt_project_dir, output_filename)

    # 獲取投影片圖片
    print(f"正在掃描 PPT 專案目錄: {ppt_project_dir}")
    image_files = get_slide_images(ppt_project_dir)

    if not image_files:
        print("錯誤：未找到任何圖片檔案")
        print("支援的格式: .png, .jpg, .jpeg, .gif, .bmp")
        print(f"\n請將投影片圖片放置在: {origin_image_dir}/")
        sys.exit(1)

    print(f"找到 {len(image_files)} 張投影片圖片\n")
    speaker_notes = load_speaker_notes(ppt_project_dir)
    if speaker_notes:
        print(f"找到 {len(speaker_notes)} 頁備註: {os.path.join(ppt_project_dir, 'speech.md')}\n")

    # 建立簡報
    print(f"正在建立 PPT (寬高比: {args.aspect_ratio})...")
    print("-" * 50)

    success = create_presentation(image_files, output_path, args.aspect_ratio, speaker_notes)

    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
